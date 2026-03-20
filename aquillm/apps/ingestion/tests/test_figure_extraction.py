"""Tests for unified document figure extraction."""

import pytest
from io import BytesIO

from aquillm.ingestion.figure_extraction.types import ExtractedFigure


class TestExtractedFigureDataclass:
    """Test the ExtractedFigure dataclass."""
    
    def test_extracted_figure_creation(self):
        figure = ExtractedFigure(
            image_bytes=b"fake-image-data",
            image_format="png",
            figure_index=0,
            nearby_text="Figure 1: Test caption",
            width=100,
            height=100,
            location_metadata={"page_number": 1},
        )
        
        assert figure.figure_index == 0
        assert figure.nearby_text == "Figure 1: Test caption"
        assert figure.location_metadata == {"page_number": 1}


class TestGenerateFigureCaption:
    """Test caption generation."""
    
    def test_caption_with_page_number(self):
        from aquillm.ingestion.figure_extraction import generate_figure_caption
        
        figure = ExtractedFigure(
            image_bytes=b"",
            image_format="png",
            figure_index=0,
            nearby_text="Figure 2: Architecture diagram",
            width=100,
            height=100,
            location_metadata={"page_number": 3},
        )
        
        caption = generate_figure_caption(figure, "Research Paper", "pdf")
        
        assert "Figure 2: Architecture diagram" in caption
        assert "Page 3" in caption
        assert "Research Paper" in caption
    
    def test_caption_with_slide_number(self):
        from aquillm.ingestion.figure_extraction import generate_figure_caption
        
        figure = ExtractedFigure(
            image_bytes=b"",
            image_format="png",
            figure_index=0,
            nearby_text="",
            width=100,
            height=100,
            location_metadata={"slide_number": 5, "slide_title": "Overview"},
        )
        
        caption = generate_figure_caption(figure, "Presentation", "pptx")
        
        assert "Slide 5" in caption
        assert "Overview" in caption
    
    def test_caption_with_sheet_name(self):
        from aquillm.ingestion.figure_extraction import generate_figure_caption
        
        figure = ExtractedFigure(
            image_bytes=b"",
            image_format="png",
            figure_index=0,
            nearby_text="",
            width=100,
            height=100,
            location_metadata={"sheet_name": "Dashboard"},
        )
        
        caption = generate_figure_caption(figure, "Spreadsheet", "xlsx")
        
        assert "Sheet: Dashboard" in caption
        assert "Spreadsheet" in caption


class TestPDFExtraction:
    """Test PDF figure extraction."""
    
    @pytest.fixture
    def fitz(self):
        return pytest.importorskip("fitz")
    
    def test_empty_pdf_returns_no_figures(self, fitz):
        from aquillm.ingestion.figure_extraction import extract_figures_from_document
        
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((100, 100), "Hello World")
        pdf_bytes = doc.tobytes()
        doc.close()
        
        figures = list(extract_figures_from_document(pdf_bytes, "pdf"))
        assert len(figures) == 0
    
    def test_pdf_with_image_extracts_figure(self, fitz):
        from aquillm.ingestion.figure_extraction import extract_figures_from_document
        from PIL import Image
        import random
        
        # Create a larger image with noise to exceed MIN_IMAGE_BYTES (5KB)
        img = Image.new('RGB', (400, 400))
        pixels = img.load()
        random.seed(42)
        for i in range(400):
            for j in range(400):
                pixels[i, j] = (random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))
        
        img_bytes = BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        
        doc = fitz.open()
        page = doc.new_page()
        page.insert_image(fitz.Rect(50, 50, 450, 450), stream=img_bytes.getvalue())
        pdf_bytes = doc.tobytes()
        doc.close()
        
        figures = list(extract_figures_from_document(pdf_bytes, "pdf"))
        
        assert len(figures) == 1
        assert figures[0].location_metadata.get("page_number") == 1
        assert figures[0].width >= 100


class TestDOCXExtraction:
    """Test DOCX figure extraction."""
    
    def test_docx_without_images(self):
        pytest.importorskip("docx")
        from aquillm.ingestion.figure_extraction import extract_figures_from_document
        from docx import Document
        
        doc = Document()
        doc.add_paragraph("Hello World")
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        
        figures = list(extract_figures_from_document(buffer.getvalue(), "docx"))
        assert len(figures) == 0


class TestPPTXExtraction:
    """Test PPTX figure extraction."""
    
    def test_pptx_without_images(self):
        pytest.importorskip("pptx")
        from aquillm.ingestion.figure_extraction import extract_figures_from_document
        from pptx import Presentation
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        figures = list(extract_figures_from_document(buffer.getvalue(), "pptx"))
        assert len(figures) == 0


class TestUnknownFormat:
    """Test handling of unknown formats."""
    
    def test_unknown_format_returns_empty(self):
        from aquillm.ingestion.figure_extraction import extract_figures_from_document
        
        figures = list(extract_figures_from_document(b"fake data", "unknown_format"))
        assert len(figures) == 0
