"""
PowerPoint PPTX presentation parsing.
"""

import io


def extract_pptx_text(data: bytes) -> str:
    """Extract text content from PPTX bytes."""
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        raise ValueError("python-pptx is required for .pptx extraction.") from exc
    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                lines.append(text.strip())
    return "\n".join(lines)


__all__ = ['extract_pptx_text']
