"""DOCX and PPTX figure extraction.

Supports two extraction methods:
1. PDF conversion + extraction - converts to PDF using LibreOffice, then extracts figures
2. Direct raster extraction - extracts embedded images directly (fallback)
"""

import io
import structlog
from typing import Iterator

from .types import ExtractedFigure

logger = structlog.stdlib.get_logger(__name__)

MIN_IMAGE_WIDTH = 100
MIN_IMAGE_HEIGHT = 100
MIN_IMAGE_BYTES = 5_000
MAX_IMAGES_PER_DOCUMENT = 50


def _get_image_dimensions(image_bytes: bytes) -> tuple[int, int]:
    """Get image dimensions using PIL."""
    try:
        from PIL import Image
        with Image.open(io.BytesIO(image_bytes)) as img:
            return img.width, img.height
    except Exception:
        return 0, 0


def _normalize_image(image_bytes: bytes, content_type: str) -> tuple[bytes, str]:
    """Normalize image to PNG if needed, return (bytes, format)."""
    ext = "png"
    if "jpeg" in content_type or "jpg" in content_type:
        ext = "jpeg"
    elif "png" in content_type:
        ext = "png"
    elif "webp" in content_type:
        ext = "webp"
    else:
        try:
            from PIL import Image
            with Image.open(io.BytesIO(image_bytes)) as img:
                if img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                output = io.BytesIO()
                img.save(output, format="PNG")
                return output.getvalue(), "png"
        except Exception:
            pass
    return image_bytes, ext


def _extract_via_pdf(data: bytes, source_format: str, filename: str) -> Iterator[ExtractedFigure] | None:
    """
    Try to extract figures by converting to PDF first.

    Returns None if conversion is not available or fails.
    """
    try:
        from .office_convert import convert_to_pdf, is_libreoffice_available
        from .pdf import extract_figures as extract_pdf_figures
    except ImportError:
        return None

    if not is_libreoffice_available():
        logger.debug("obs.ingest.libreoffice_unavailable", source_format=source_format)
        return None

    pdf_bytes = convert_to_pdf(data, source_format, filename)
    if not pdf_bytes:
        logger.debug("obs.ingest.pdf_conversion_failed", filename=filename or source_format)
        return None

    logger.info("obs.ingest.pdf_conversion_used", source_format=source_format.upper(), filename=filename)

    # Update location metadata to indicate conversion was used
    for figure in extract_pdf_figures(pdf_bytes, filename):
        figure.location_metadata["converted_from"] = source_format
        yield figure


def _extract_docx_direct(data: bytes, filename: str = "") -> Iterator[ExtractedFigure]:
    """Extract embedded raster images directly from DOCX."""
    try:
        from docx import Document
    except ImportError:
        logger.warning("obs.ingest.docx_dependency_missing", dependency="python-docx")
        return

    try:
        doc = Document(io.BytesIO(data))
    except Exception as exc:
        logger.warning("obs.ingest.docx_open_failed", error_type=type(exc).__name__, error=str(exc))
        return

    total_extracted = 0

    try:
        for rel in doc.part.rels.values():
            if total_extracted >= MAX_IMAGES_PER_DOCUMENT:
                break

            if "image" not in rel.reltype:
                continue

            try:
                image_part = rel.target_part
                image_bytes = image_part.blob

                if not image_bytes or len(image_bytes) < MIN_IMAGE_BYTES:
                    continue

                width, height = _get_image_dimensions(image_bytes)
                if width < MIN_IMAGE_WIDTH or height < MIN_IMAGE_HEIGHT:
                    continue

                content_type = getattr(image_part, 'content_type', 'image/png')
                image_bytes, img_format = _normalize_image(image_bytes, content_type)

                yield ExtractedFigure(
                    image_bytes=image_bytes,
                    image_format=img_format,
                    figure_index=total_extracted,
                    nearby_text="",
                    width=width,
                    height=height,
                    location_metadata={"source": "docx_embedded", "extraction_method": "direct"},
                )

                total_extracted += 1

            except Exception as exc:
                logger.debug("obs.ingest.docx_image_extract_failed", error_type=type(exc).__name__, error=str(exc))
                continue

    except Exception as exc:
        logger.warning("obs.ingest.docx_extraction_failed", error_type=type(exc).__name__, error=str(exc))

    if total_extracted > 0:
        logger.info("obs.ingest.figures_docx_done", filename=filename, figure_count=total_extracted)


def _extract_pptx_direct(data: bytes, filename: str = "") -> Iterator[ExtractedFigure]:
    """Extract embedded raster images directly from PPTX."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.warning("obs.ingest.pptx_dependency_missing", dependency="python-pptx")
        return

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        logger.warning("obs.ingest.pptx_open_failed", error_type=type(exc).__name__, error=str(exc))
        return

    total_extracted = 0

    try:
        for slide_num, slide in enumerate(prs.slides, start=1):
            if total_extracted >= MAX_IMAGES_PER_DOCUMENT:
                break

            slide_title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        slide_title = shape.text.strip()[:200]
                        break

            for shape in slide.shapes:
                if total_extracted >= MAX_IMAGES_PER_DOCUMENT:
                    break

                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue

                try:
                    image = shape.image
                    image_bytes = image.blob

                    if not image_bytes or len(image_bytes) < MIN_IMAGE_BYTES:
                        continue

                    width, height = _get_image_dimensions(image_bytes)
                    if width < MIN_IMAGE_WIDTH or height < MIN_IMAGE_HEIGHT:
                        continue

                    content_type = image.content_type or 'image/png'
                    image_bytes, img_format = _normalize_image(image_bytes, content_type)

                    alt_text = ""
                    try:
                        if hasattr(shape, '_element'):
                            desc_elem = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
                            if desc_elem is not None:
                                alt_text = desc_elem.get('descr', '')[:500]
                    except Exception:
                        pass

                    nearby_text = alt_text or slide_title

                    yield ExtractedFigure(
                        image_bytes=image_bytes,
                        image_format=img_format,
                        figure_index=total_extracted,
                        nearby_text=nearby_text,
                        width=width,
                        height=height,
                        location_metadata={
                            "slide_number": slide_num,
                            "slide_title": slide_title,
                            "extraction_method": "direct",
                        },
                    )

                    total_extracted += 1

                except Exception as exc:
                    logger.debug("obs.ingest.pptx_image_extract_failed", slide_num=slide_num, error_type=type(exc).__name__, error=str(exc))
                    continue

    except Exception as exc:
        logger.warning("obs.ingest.pptx_extraction_failed", error_type=type(exc).__name__, error=str(exc))

    if total_extracted > 0:
        logger.info("obs.ingest.figures_pptx_done", filename=filename, figure_count=total_extracted)


def extract_figures_docx(data: bytes, filename: str = "") -> Iterator[ExtractedFigure]:
    """
    Extract figures from a DOCX file.

    First tries PDF conversion (to capture vector graphics like charts),
    falls back to direct extraction if LibreOffice is not available.

    Args:
        data: Raw DOCX bytes
        filename: Optional filename for logging

    Yields:
        ExtractedFigure for each valid image found
    """
    # Try PDF conversion first (captures vector graphics)
    pdf_result = _extract_via_pdf(data, "docx", filename)
    if pdf_result is not None:
        yield from pdf_result
        return

    # Fall back to direct extraction
    logger.debug("obs.ingest.docx_fallback_direct", filename=filename)
    yield from _extract_docx_direct(data, filename)


def extract_figures_pptx(data: bytes, filename: str = "") -> Iterator[ExtractedFigure]:
    """
    Extract figures from a PPTX file.

    First tries PDF conversion (to capture vector graphics like charts and diagrams),
    falls back to direct extraction if LibreOffice is not available.

    Args:
        data: Raw PPTX bytes
        filename: Optional filename for logging

    Yields:
        ExtractedFigure for each valid image found
    """
    # Try PDF conversion first (captures vector graphics)
    pdf_result = _extract_via_pdf(data, "pptx", filename)
    if pdf_result is not None:
        yield from pdf_result
        return

    # Fall back to direct extraction
    logger.debug("obs.ingest.pptx_fallback_direct", filename=filename)
    yield from _extract_pptx_direct(data, filename)
